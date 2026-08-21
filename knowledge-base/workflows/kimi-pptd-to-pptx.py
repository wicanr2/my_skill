"""PPTD → PPTX 轉換器(python-pptx)。搭配 kimi-pptd-deck-authoring.md 使用。

為什麼不用 open-kimi-ppt 內建的 export_pptx.py:那支走 Kimi 公開編輯器的
瀏覽器端 writer,實測三個階段都會失敗,而且**用 skill 自帶的 example/yu7-ppt
做正對照也一樣失敗** —— 是環境/上游問題。改用 python-pptx 在本機直接寫 OOXML,
確定性、無遠端依賴。

涵蓋 text(含 <p>/<span style>/<strong> 富文字)、shape(rect)、line(含箭頭與
虛線,折線拆成多段直線)、table(含合併與逐邊框線)。
**不含 image 與 chart** —— 需要圖表就用 shape 自己畫長條(遠端 chart 渲染會
整張消失,見筆記)。要支援 image 再自行補 slide.shapes.add_picture。

用法:
    python3 kimi-pptd-to-pptx.py <deck.pptd> <out.pptx>

PPTD 規格:1px = 1pt,原點左上。960×540pt 正好是 16:9 寬螢幕(13.333×7.5 in)。

[HARD] 一律 docker + uv 跑,不污染系統 Python:
    docker run --rm -v "$(pwd)":/work -w /work \
      ghcr.io/astral-sh/uv:python3.12-bookworm-slim \
      bash -c "uv run --with python-pptx --with pyyaml kimi-pptd-to-pptx.py deck.pptd deck.pptx"
"""
import copy
import html
import os
import re
import sys
from html.parser import HTMLParser

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu

SRC = sys.argv[1] if len(sys.argv) > 1 else "deck.pptd"
OUT = sys.argv[2] if len(sys.argv) > 2 else "deck.pptx"
BASE = os.path.dirname(os.path.abspath(SRC)) or "."

# 字型:PPTX 交付對象多半是 Windows PowerPoint,中文用微軟正黑體最穩;
# 西文用 Arial(PPTD 原設計的 Liter/MiSans 幾乎不會裝在對方機器上)。
FONT_EA = "Microsoft JhengHei"
FONT_LATIN = "Arial"

deck = yaml.safe_load(open(SRC, encoding="utf-8"))
PW, PH = deck["size"]
theme = deck.get("theme") or {}
COLORS = theme.get("colors") or {}
TSTYLES = theme.get("textStyles") or {}
TBSTYLES = theme.get("tableStyles") or {}

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
          "bottom": MSO_ANCHOR.BOTTOM}


def rgb(c, default="000000"):
    """解析 PPTD Color:$主題引用 / #RRGGBB / #RRGGBBAA(捨去 alpha)。"""
    if not c:
        return RGBColor.from_string(default)
    c = str(c)
    if c.startswith("$"):
        c = COLORS.get(c[1:], "#" + default)
    c = c.lstrip("#")
    return RGBColor.from_string(c[:6].upper())


def style_of(name):
    return dict(TSTYLES.get(str(name).lstrip("$"), {})) if name else {}


# ── 富文字解析 ─────────────────────────────────────────────────────────────
class RichParser(HTMLParser):
    """把 <p>/<span style>/<strong>/<em> 解析成 [(段落屬性, [(文字, 行內屬性)])]。

    只支援本簡報用到的標籤。遇到沒實作的標籤就當透明(不丟棄內容)——
    寧可少一層樣式,也不要吃掉文字。
    """

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.paras = []
        self.cur = None
        self.stack = []

    def _ensure(self, pattr=None):
        if self.cur is None:
            self.paras.append((pattr or {}, []))
            self.cur = self.paras[-1][1]

    def handle_starttag(self, tag, attrs):
        a = dict(attrs)
        if tag == "p":
            self.paras.append((self._pstyle(a.get("style", "")), []))
            self.cur = self.paras[-1][1]
        elif tag == "span":
            self.stack.append(self._rstyle(a.get("style", "")))
        elif tag in ("strong", "b"):
            self.stack.append({"bold": True})
        elif tag in ("em", "i"):
            self.stack.append({"italic": True})
        elif tag == "br":
            self._ensure()
            self.cur.append(("\n", self._merged()))

    def handle_endtag(self, tag):
        if tag == "p":
            self.cur = None
        elif tag in ("span", "strong", "b", "em", "i"):
            if self.stack:
                self.stack.pop()

    def handle_data(self, data):
        if not data:
            return
        self._ensure()
        self.cur.append((data, self._merged()))

    def _merged(self):
        m = {}
        for s in self.stack:
            m.update(s)
        return m

    @staticmethod
    def _pstyle(s):
        out = {}
        for k, v in re.findall(r"([a-z-]+)\s*:\s*([^;]+)", s or ""):
            if k == "text-align":
                out["align"] = v.strip()
        return out

    @staticmethod
    def _rstyle(s):
        out = {}
        for k, v in re.findall(r"([a-z-]+)\s*:\s*([^;]+)", s or ""):
            v = v.strip()
            if k == "color":
                out["color"] = v
            elif k == "font-size":
                out["fontSize"] = float(re.sub(r"[^\d.]", "", v) or 0)
            elif k == "font-family":
                out["fontFamily"] = v
            elif k == "font-weight":
                out["bold"] = v in ("bold", "700", "800", "900")
        return out


def parse_rich(text):
    t = str(text)
    if "<" not in t:
        # 純文字:\n 當換行
        return [({}, [(line, {})]) for line in t.split("\n")]
    p = RichParser()
    p.feed(t)
    p.close()
    return p.paras or [({}, [(re.sub(r"<[^>]+>", "", t), {})])]


def apply_run(run, base, over):
    st = dict(base)
    st.update(over or {})
    f = run.font
    f.size = Pt(float(st.get("fontSize", 18)))
    f.bold = bool(st.get("bold", False))
    f.italic = bool(st.get("italic", False))
    f.color.rgb = rgb(st.get("color"))
    latin = FONT_LATIN
    ff = st.get("fontFamily")
    if isinstance(ff, dict):
        latin = ff.get("latin", FONT_LATIN)
    f.name = latin
    # 東亞字型要另外寫進 rPr,python-pptx 的 font.name 只設 latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT_EA)
    sp = st.get("letterSpacing")
    if sp:
        rPr.set("spc", str(int(float(sp) * 100)))


def fill_textframe(tf, content, wrap_default=True):
    base = style_of(content.get("style"))
    for k in ("color", "fontSize", "fontFamily", "bold", "italic",
              "lineHeight", "letterSpacing"):
        if k in content:
            base[k] = content[k]
    base.setdefault("fontSize", 18)

    align = content.get("align") or ["left", "top"]
    tf.word_wrap = bool(content.get("wrap", wrap_default))
    tf.vertical_anchor = ANCHOR.get(align[1], MSO_ANCHOR.TOP)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = parse_rich(content.get("text", ""))
    for i, (pattr, runs) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN.get(pattr.get("align", align[0]), PP_ALIGN.LEFT)
        lh = base.get("lineHeight")
        if lh:
            p.line_spacing = float(lh)
        if not runs:
            runs = [("", {})]
        for txt, over in runs:
            r = p.add_run()
            r.text = html.unescape(txt)
            apply_run(r, base, over)


def set_line(shape_line, border, default_none=True):
    if not border:
        if default_none:
            shape_line.fill.background()
        return
    shape_line.color.rgb = rgb(border.get("color"))
    shape_line.width = Pt(float(border.get("width", 1)))
    st = border.get("style", "solid")
    ln = shape_line._get_or_add_ln()
    for d in ln.findall(qn("a:prstDash")):
        ln.remove(d)
    if st in ("dash", "dot"):
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash" if st == "dash" else "sysDot"})
        ln.append(d)


def add_text(slide, e):
    b = e["bounds"]
    box = slide.shapes.add_textbox(Pt(b[0]), Pt(b[1]), Pt(b[2]), Pt(b[3]))
    fill_textframe(box.text_frame, e["content"])
    return box


def add_shape(slide, e):
    b = e["bounds"]
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(b[0]), Pt(b[1]),
                                Pt(b[2]), Pt(b[3]))
    f = e.get("fill")
    if f and f.get("type") == "solid":
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(f.get("color"))
    else:
        sp.fill.background()
    set_line(sp.line, e.get("border"))
    sp.shadow.inherit = False
    if sp.has_text_frame:                      # PPTD 的 shape 不帶文字
        sp.text_frame.text = ""
    return sp


def add_line(slide, e):
    """PPTD line 用 viewBox + points;本簡報只有直線與直角折線。

    直線直接畫連接線;折線(如迴圈那條)拆成多段直線畫,
    箭頭只加在最後一段 —— 這樣不需要自由造型幾何,PowerPoint 相容性最好。
    """
    b = e["bounds"]
    vb = e.get("viewBox") or [max(b[2], 1), max(b[3], 1)]
    sx = b[2] / (vb[0] or 1)
    sy = b[3] / (vb[1] or 1)
    pts = []
    for p in str(e.get("points", "")).split():
        try:
            px, py = (float(t) for t in p.split(","))
        except ValueError:
            continue
        pts.append((b[0] + px * sx, b[1] + py * sy))
    if len(pts) < 2:
        return None
    arrow = (e.get("arrow") or [None, None])[1]
    out = None
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Pt(x1), Pt(y1), Pt(x2), Pt(y2))
        set_line(cn.line, e.get("border") or {"width": 1, "color": "#000000"},
                 default_none=False)
        if arrow and i == len(pts) - 2:
            ln = cn.line._get_or_add_ln()
            tail = ln.makeelement(qn("a:tailEnd"),
                                  {"type": "stealth" if arrow == "stealth" else "triangle",
                                   "w": "med", "len": "med"})
            ln.append(tail)
        out = cn
    return out


def cell_border(cell, spec):
    """PPTD BorderSpec → OOXML 的 lnT/lnR/lnB/lnL。python-pptx 沒有這個 API。"""
    if spec is None:
        sides = [None] * 4
    elif isinstance(spec, dict):
        sides = [spec] * 4
    elif isinstance(spec, list) and len(spec) == 2:
        sides = [spec[0], spec[1], spec[0], spec[1]]
    elif isinstance(spec, list) and len(spec) == 4:
        sides = spec
    else:
        sides = [None] * 4
    tcPr = cell._tc.get_or_add_tcPr()
    for tag, bd in zip(("a:lnT", "a:lnR", "a:lnB", "a:lnL"), sides):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        ln = tcPr.makeelement(qn(tag), {})
        if bd:
            ln.set("w", str(int(float(bd.get("width", 1)) * 12700)))
            ln.set("cap", "flat")
            fill = ln.makeelement(qn("a:solidFill"), {})
            clr = ln.makeelement(qn("a:srgbClr"), {"val": str(rgb(bd.get("color")))})
            fill.append(clr)
            ln.append(fill)
        else:
            ln.append(ln.makeelement(qn("a:noFill"), {}))
        # OOXML 要求 lnL/lnR/lnT/lnB 依序排在 tcPr 開頭
        tcPr.insert(0, ln)


def add_table(slide, e):
    b = e["bounds"]
    rows = e["rows"]
    colw, rowh = e["columnWidths"], e["rowHeights"]
    ncol = len(colw)
    gf = slide.shapes.add_table(len(rows), ncol, Pt(b[0]), Pt(b[1]),
                                Pt(b[2]), Pt(b[3]))
    tbl = gf.table
    # 關掉 PowerPoint 預設表格樣式的首列強調與斑馬紋 —— 本簡報自己畫線
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for a in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
            tblPr.set(a, "0")
        for st in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(st)
    for i, w in enumerate(colw):
        tbl.columns[i].width = Emu(int(Pt(b[2]).emu * w))
    for i, h in enumerate(rowh):
        tbl.rows[i].height = Emu(int(Pt(b[3]).emu * h))

    style = e.get("style")
    ts = TBSTYLES.get(str(style).lstrip("$"), {}) if isinstance(style, str) else (style or {})
    cs_base = ts.get("cellStyle") or {}
    first = ts.get("firstRowStyle") or {}
    last = ts.get("lastRowStyle") or {}

    occ = [[False] * ncol for _ in rows]
    for ri, row in enumerate(rows):
        ci = 0
        for cd in row:
            while ci < ncol and occ[ri][ci]:
                ci += 1
            if ci >= ncol:
                break
            rs, cs = cd.get("rowSpan", 1), cd.get("colSpan", 1)
            for dr in range(rs):
                for dc in range(cs):
                    if ri + dr < len(rows) and ci + dc < ncol:
                        occ[ri + dr][ci + dc] = True
            cell = tbl.cell(ri, ci)
            if rs > 1 or cs > 1:
                cell.merge(tbl.cell(min(ri + rs - 1, len(rows) - 1),
                                    min(ci + cs - 1, ncol - 1)))
            st = dict(cs_base)
            if ri == 0 and first:
                st.update(first)
            if ri == len(rows) - 1 and last:
                st.update(last)
            st.update({k: v for k, v in cd.items()
                       if k in ("color", "fontSize", "bold", "italic",
                                "fontFamily", "lineHeight", "align")})
            f = cd.get("fill") or st.get("fill")
            if f and f.get("type") == "solid":
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(f.get("color"))
            else:
                cell.fill.background()
            bd = cd.get("border", st.get("border"))
            cell_border(cell, bd)
            cell.margin_left = cell.margin_right = Pt(5)
            cell.margin_top = cell.margin_bottom = Pt(2)
            al = st.get("align") or ["left", "middle"]
            cell.vertical_anchor = ANCHOR.get(al[1], MSO_ANCHOR.MIDDLE)
            content = {"text": cd.get("text", ""), "align": al}
            for k in ("color", "fontSize", "bold", "italic", "lineHeight"):
                if k in st:
                    content[k] = st[k]
            if cd.get("textStyle"):
                content["style"] = cd["textStyle"]
            fill_textframe(cell.text_frame, content)
            ci += cs
    return gf


def page_background(slide, bg):
    if not bg or bg.get("type") != "solid":
        return
    cSld = slide._element.find(qn("p:cSld"))
    bgEl = cSld.makeelement(qn("p:bg"), {})
    bgPr = cSld.makeelement(qn("p:bgPr"), {})
    fill = cSld.makeelement(qn("a:solidFill"), {})
    fill.append(cSld.makeelement(qn("a:srgbClr"), {"val": str(rgb(bg.get("color")))}))
    bgPr.append(fill)
    bgPr.append(cSld.makeelement(qn("a:effectLst"), {}))
    bgEl.append(bgPr)
    cSld.insert(0, bgEl)


def add_fade(slide):
    """每張投影片加淡入淡出切換。

    [HARD] transition 必須是 <p:sld> 的直接子元素,而且要排在 cSld / clrMapOvr 之後。
    塞進 cSld 裡面 Office 會直接忽略(用 grep <p:fade> 驗不出這個差別)。
    """
    sld = slide._element
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    tr = sld.makeelement(qn("p:transition"), {"spd": "med"})
    tr.append(sld.makeelement(qn("p:fade"), {}))
    clr = sld.find(qn("p:clrMapOvr"))
    if clr is not None:
        clr.addnext(tr)
    else:
        cSld = sld.find(qn("p:cSld"))
        cSld.addnext(tr)


prs = Presentation()
prs.slide_width = Pt(PW)
prs.slide_height = Pt(PH)
blank = prs.slide_layouts[6]

HANDLERS = {"text": add_text, "shape": add_shape, "line": add_line,
            "table": add_table}
skipped = []
for rel in deck["pages"]:
    pg = yaml.safe_load(open(os.path.join(BASE, rel), encoding="utf-8"))
    slide = prs.slides.add_slide(blank)
    page_background(slide, pg.get("background") or {"type": "solid", "color": "#FFFFFF"})
    for e in pg.get("elements") or []:
        h = HANDLERS.get(e["elementType"])
        if h is None:
            skipped.append(f"{rel}:{e.get('elementId')} ({e['elementType']})")
            continue
        h(slide, e)
    if pg.get("notes"):
        slide.notes_slide.notes_text_frame.text = pg["notes"]
    add_fade(slide)

prs.save(OUT)
print(f"寫出 {OUT}:{len(prs.slides.__iter__.__self__._sldIdLst)} 頁, "
      f"{os.path.getsize(OUT)} bytes")
if skipped:
    print("⚠ 未處理的元素:")
    for s in skipped:
        print("   ", s)
